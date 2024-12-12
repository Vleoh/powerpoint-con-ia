# Lógica para manejar la inserción de contenido en las plantillas

from pptx import Presentation
from pptx.util import Inches, Pt

class TemplateManager:
    def __init__(self):
        self.presentation = None
    
    def create_presentation(self):
        """Creates a new presentation"""
        self.presentation = Presentation()
        # Define a file path to save the presentation
        filename = 'presentation.pptx'
        self.save_presentation(filename)
        return filename
    
    def add_title_slide(self, title, subtitle=None):
        """Adds a title slide to the presentation"""
        layout = self.presentation.slide_layouts[0]
        slide = self.presentation.slides.add_slide(layout)
        
        title_placeholder = slide.shapes.title
        subtitle_placeholder = slide.placeholders[1]
        
        title_placeholder.text = title
        if subtitle:
            subtitle_placeholder.text = subtitle
    
    def add_content_slide(self, title, content):
        """Adds a content slide to the presentation"""
        layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(layout)
        
        title_placeholder = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        
        title_placeholder.text = title
        content_placeholder.text = content
    
    def save_presentation(self, filename):
        """Saves the presentation to a file"""
        self.presentation.save(filename)

def insert_content(template, content):
    # Aquí irá el código para insertar contenido en la plantilla
    pass
